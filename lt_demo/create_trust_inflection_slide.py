from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
from pathlib import Path

_SCRIPT_DIR = Path(__file__).resolve().parent
screenshot_folder = str(_SCRIPT_DIR.parent / 'screenshots')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(15, 15, 15)
bg.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The Trust Inflection Point in AI Restyling"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.333), Inches(0.4))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "From 'Obviously AI' to 'I'd Actually Share This'"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(150, 150, 150)
p.alignment = PP_ALIGN.CENTER

# Row labels
before_label = slide.shapes.add_textbox(Inches(0.1), Inches(1.8), Inches(1), Inches(0.5))
p = before_label.text_frame.paragraphs[0]
p.text = "Pre-\nInflection"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 100, 100)
p.alignment = PP_ALIGN.CENTER

after_label = slide.shapes.add_textbox(Inches(0.1), Inches(4.4), Inches(1), Inches(0.5))
p = after_label.text_frame.paragraphs[0]
p.text = "Post-\nInflection"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 255, 100)
p.alignment = PP_ALIGN.CENTER

# Categories and images - using our actual AI Restyle outputs
# Before = more stylized/cartoonish, After = more subtle/artistic
pairs = [
    ("Portrait", "img1-04-chibi-sticker.png", "img1-11-pencil-portrait.png"),
    ("Creative", "img1-05-caricature.png", "img1-03-anime.png"),
    ("Artistic", "img1-01-movie-poster.png", "img1-12-storybook.png"),
    ("Style", "img1-08-graffiti.png", "img1-14-pop-art.png"),
]

for i, (category, before_img, after_img) in enumerate(pairs):
    x = 1.2 + i * 3.0

    # Category label
    cat_box = slide.shapes.add_textbox(Inches(x), Inches(1.05), Inches(2.8), Inches(0.35))
    p = cat_box.text_frame.paragraphs[0]
    p.text = category
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Before image
    before_path = os.path.join(screenshot_folder, before_img)
    if os.path.exists(before_path):
        # Red border box
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x-0.05), Inches(1.4), Inches(2.9), Inches(2.55))
        border.fill.background()
        border.line.color.rgb = RGBColor(180, 60, 60)
        border.line.width = Pt(2)

        slide.shapes.add_picture(before_path, Inches(x), Inches(1.45), width=Inches(2.8))

    # After image
    after_path = os.path.join(screenshot_folder, after_img)
    if os.path.exists(after_path):
        # Green border box
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x-0.05), Inches(4.05), Inches(2.9), Inches(2.55))
        border.fill.background()
        border.line.color.rgb = RGBColor(60, 180, 60)
        border.line.width = Pt(2)

        slide.shapes.add_picture(after_path, Inches(x), Inches(4.1), width=Inches(2.8))

# Arrow between rows
arrow_box = slide.shapes.add_textbox(Inches(6.2), Inches(3.95), Inches(1), Inches(0.3))
p = arrow_box.text_frame.paragraphs[0]
p.text = "▼"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 200, 255)
p.alignment = PP_ALIGN.CENTER

# Bottom insight
insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
tf = insight_box.text_frame
p = tf.paragraphs[0]
p.text = "The inflection point: When AI editing becomes subtle enough to trust and share."
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# Save
output_path = str(_SCRIPT_DIR / 'Trust_Inflection_Point_Slide.pptx')
prs.save(output_path)
print(f'Slide saved to: {output_path}')
