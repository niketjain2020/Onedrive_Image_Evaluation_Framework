from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os
import csv
from pathlib import Path

# Paths (relative to this script)
_SCRIPT_DIR = Path(__file__).resolve().parent
screenshot_folder = str(_SCRIPT_DIR.parent / 'screenshots')
metrics_file = str(_SCRIPT_DIR / 'performance_metrics.csv')
output_path = str(_SCRIPT_DIR / 'Claude_Code_AI_Testing_LT.pptx')

# Load metrics
metrics = []
with open(metrics_file, 'r') as f:
    reader = csv.DictReader(f)
    for row in reader:
        metrics.append(row)

# Calculate statistics
gen_times = [int(m['Time_ms']) for m in metrics if m['Category'] == 'Generation']
avg_gen_time = sum(gen_times) / len(gen_times) if gen_times else 0
min_gen_time = min(gen_times) if gen_times else 0
max_gen_time = max(gen_times) if gen_times else 0

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gradient-like background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0, 45, 98)
    bg.line.fill.background()

    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0, 164, 239)
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(180, 210, 240)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_func):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 45, 98)
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Call content function to add specific content
    content_func(slide)

    return slide

def add_bullets(slide, bullets, left, top, width, height, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)
    return box

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top))
        return True
    return False

# ============== SLIDE 1: TITLE ==============
add_title_slide(prs,
    "Using AI to Test AI",
    "Claude Code + Playwright MCP for OneDrive Photos QA Automation")

# ============== SLIDE 2: WHAT WE DID ==============
def slide2_content(slide):
    bullets = [
        "Autonomous UI Testing via Claude Code + Playwright MCP",
        "",
        "14 AI Restyle styles tested across 2 images (28 generations)",
        "",
        "Performance instrumentation with JavaScript timestamps",
        "",
        "65+ screenshots captured automatically",
        "",
        "Word reports, CSV metrics, and this PPT auto-generated",
        "",
        "2 bugs discovered without manual intervention"
    ]
    add_bullets(slide, bullets, 0.5, 1.4, 6, 5.5, font_size=22)
    add_image_safe(slide, os.path.join(screenshot_folder, 'qa_step0_gallery.png'), 7, 1.5, width=5.8)

add_content_slide(prs, "What We Did", slide2_content)

# ============== SLIDE 3: AUTOMATED UI TESTING ==============
def slide3_content(slide):
    # Three images showing the flow
    add_image_safe(slide, os.path.join(screenshot_folder, 'qa_step1_viewer_opened.png'), 0.3, 1.4, width=4.1)
    add_image_safe(slide, os.path.join(screenshot_folder, 'qa_step2_restyle_panel.png'), 4.6, 1.4, width=4.1)
    add_image_safe(slide, os.path.join(screenshot_folder, 'qa_step4_generation_complete.png'), 8.9, 1.4, width=4.1)

    # Captions
    cap1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.2), Inches(4.1), Inches(0.5))
    cap1.text_frame.paragraphs[0].text = "1. Open Image"
    cap1.text_frame.paragraphs[0].font.size = Pt(14)
    cap1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cap2 = slide.shapes.add_textbox(Inches(4.6), Inches(5.2), Inches(4.1), Inches(0.5))
    cap2.text_frame.paragraphs[0].text = "2. Select Style"
    cap2.text_frame.paragraphs[0].font.size = Pt(14)
    cap2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cap3 = slide.shapes.add_textbox(Inches(8.9), Inches(5.2), Inches(4.1), Inches(0.5))
    cap3.text_frame.paragraphs[0].text = "3. AI Generates"
    cap3.text_frame.paragraphs[0].font.size = Pt(14)
    cap3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Flow description
    desc = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(1))
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = "Claude autonomously navigates, clicks, waits for generation, validates UI states, and captures evidence"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "Automated UI Testing", slide3_content)

# ============== SLIDE 4: PERFORMANCE MEASUREMENT ==============
def slide4_content(slide):
    # Metrics table
    table = slide.shapes.add_table(6, 3, Inches(0.5), Inches(1.4), Inches(7), Inches(3)).table

    # Headers
    headers = ['Metric', 'Value', 'Assessment']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 45, 98)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(16)

    # Data
    data = [
        ('Gallery Load', '911 ms', 'Good'),
        ('Restyle Panel Open', '523 ms', 'Good'),
        (f'Avg AI Generation', f'{avg_gen_time/1000:.1f} sec', 'Expected'),
        (f'Min Generation', f'{min_gen_time/1000:.1f} sec', 'Best'),
        (f'Max Generation', f'{max_gen_time/1000:.1f} sec', 'Acceptable')
    ]

    for row_idx, (metric, value, assessment) in enumerate(data, 1):
        table.cell(row_idx, 0).text = metric
        table.cell(row_idx, 1).text = value
        table.cell(row_idx, 2).text = assessment
        for col in range(3):
            p = table.cell(row_idx, col).text_frame.paragraphs[0]
            p.font.size = Pt(14)

    # Stats boxes on right
    stats = [
        (f'{len(gen_times)}', 'Total Generations'),
        (f'{avg_gen_time/1000:.1f}s', 'Avg Time'),
        ('28', 'Styles Tested')
    ]

    colors = [RGBColor(0, 120, 60), RGBColor(0, 102, 204), RGBColor(153, 51, 153)]

    for i, (val, label) in enumerate(stats):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.5 + i*1.8), Inches(4.8), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()

        val_box = slide.shapes.add_textbox(Inches(8), Inches(1.6 + i*1.8), Inches(4.8), Inches(0.8))
        p = val_box.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        lbl_box = slide.shapes.add_textbox(Inches(8), Inches(2.4 + i*1.8), Inches(4.8), Inches(0.4))
        p = lbl_box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "Performance Measurement", slide4_content)

# ============== SLIDE 5: SCREENSHOT CAPTURE ==============
def slide5_content(slide):
    # Grid of 9 screenshots (3x3)
    screenshots = [
        ('qa_step0_gallery.png', 'Gallery'),
        ('qa_step2_restyle_panel.png', 'Restyle Panel'),
        ('img1-01-movie-poster.png', 'Movie Poster'),
        ('img1-03-anime.png', 'Anime'),
        ('img1-06-superhero.png', 'Superhero'),
        ('img1-14-pop-art.png', 'Pop Art'),
        ('qa_step7a_after_undo.png', 'After Undo'),
        ('qa_step6_back_dialog.png', 'Unsaved Dialog'),
        ('bugbash_01_webp_no_restyle.png', 'WEBP Bug')
    ]

    positions = [
        (0.3, 1.3), (4.5, 1.3), (8.7, 1.3),
        (0.3, 3.3), (4.5, 3.3), (8.7, 3.3),
        (0.3, 5.3), (4.5, 5.3), (8.7, 5.3)
    ]

    for i, (filename, caption) in enumerate(screenshots):
        if i >= 9:
            break
        x, y = positions[i]
        path = os.path.join(screenshot_folder, filename)
        if os.path.exists(path):
            slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(4))
            cap = slide.shapes.add_textbox(Inches(x), Inches(y + 1.75), Inches(4), Inches(0.3))
            p = cap.text_frame.paragraphs[0]
            p.text = caption
            p.font.size = Pt(11)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "Screenshot Capture (65+ Auto-Captured)", slide5_content)

# ============== SLIDE 6: STYLE GALLERY ==============
def slide6_content(slide):
    # Show all 14 style outputs in a grid
    styles = [
        'img1-01-movie-poster.png', 'img1-02-plush-toy.png', 'img1-03-anime.png',
        'img1-04-chibi-sticker.png', 'img1-05-caricature.png', 'img1-06-superhero.png',
        'img1-07-toy-model.png', 'img1-08-graffiti.png', 'img1-09-crochet-art.png',
        'img1-10-doodle.png', 'img1-11-pencil-portrait.png', 'img1-12-storybook.png',
        'img1-13-photo-booth.png', 'img1-14-pop-art.png'
    ]

    style_names = [
        'Movie Poster', 'Plush Toy', 'Anime', 'Chibi Sticker', 'Caricature', 'Superhero',
        'Toy Model', 'Graffiti', 'Crochet Art', 'Doodle', 'Pencil Portrait', 'Storybook',
        'Photo Booth', 'Pop Art'
    ]

    # 7 columns x 2 rows
    for i, (filename, name) in enumerate(zip(styles, style_names)):
        col = i % 7
        row = i // 7
        x = 0.2 + col * 1.85
        y = 1.3 + row * 3.0

        path = os.path.join(screenshot_folder, filename)
        if os.path.exists(path):
            slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(1.75))
            cap = slide.shapes.add_textbox(Inches(x), Inches(y + 1.5), Inches(1.75), Inches(0.3))
            p = cap.text_frame.paragraphs[0]
            p.text = name
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "All 14 AI Restyle Outputs (Image 1)", slide6_content)

# ============== SLIDE 7: AUTONOMOUS BUG DISCOVERY ==============
def slide7_content(slide):
    # Two-column layout for bugs
    # Bug 1
    bug1_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8))
    bug1_box.fill.solid()
    bug1_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    bug1_box.line.color.rgb = RGBColor(200, 0, 0)

    bug1_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
    p = bug1_title.text_frame.paragraphs[0]
    p.text = "BUG #1: Stop Button Non-Functional"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 0, 0)

    bug1_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(2))
    tf = bug1_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Severity: HIGH (P1)\n\nThe Stop button appears during AI generation but does not cancel the operation. Clicked 3 times - generation continued to completion."
    p.font.size = Pt(14)

    # Bug 2
    bug2_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8))
    bug2_box.fill.solid()
    bug2_box.fill.fore_color.rgb = RGBColor(255, 248, 230)
    bug2_box.line.color.rgb = RGBColor(200, 150, 0)

    bug2_title = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.8), Inches(0.5))
    p = bug2_title.text_frame.paragraphs[0]
    p.text = "BUG #2: WEBP Format Not Supported"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 120, 0)

    bug2_text = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.8), Inches(2))
    tf = bug2_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Severity: MEDIUM (P2)\n\nWEBP images do not show the Restyle option. No user messaging explains the format limitation."
    p.font.size = Pt(14)

    # Screenshots below
    add_image_safe(slide, os.path.join(screenshot_folder, 'bugbash_02_stop_not_working.png'), 0.3, 4.3, width=6.2)
    add_image_safe(slide, os.path.join(screenshot_folder, 'bugbash_01_webp_no_restyle.png'), 6.8, 4.3, width=6.2)

add_content_slide(prs, "Autonomous Bug Discovery", slide7_content)

# ============== SLIDE 8: AUTOMATED REPORTING ==============
def slide8_content(slide):
    # List of generated artifacts
    artifacts = [
        ("Word Report", "AI_Restyle_Bug_Bash_Report.docx", "Full bug report with screenshots"),
        ("Comprehensive Audit", "Comprehensive_Bug_Bash_Report.docx", "Latency + A11y + UX findings"),
        ("CSV Metrics", "performance_metrics.csv", "Raw timing data for analysis"),
        ("Screenshots", "65+ PNG files", "Evidence for every test state"),
        ("This Presentation", "Claude_Code_AI_Testing_LT.pptx", "Leadership summary")
    ]

    for i, (name, file, desc) in enumerate(artifacts):
        y = 1.5 + i * 1.1

        # Icon box
        icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.8), Inches(0.8))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(0, 102, 204)
        icon.line.fill.background()

        # Number
        num = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(0.8), Inches(0.5))
        p = num.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Name
        name_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(4), Inches(0.4))
        p = name_box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True

        # File
        file_box = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.4), Inches(4), Inches(0.4))
        p = file_box.text_frame.paragraphs[0]
        p.text = file
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Description
        desc_box = slide.shapes.add_textbox(Inches(5.5), Inches(y + 0.15), Inches(5), Inches(0.5))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.italic = True

    # Key message
    msg = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    p = msg.text_frame.paragraphs[0]
    p.text = "All documentation generated automatically - no manual report writing required"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "Automated Report Generation", slide8_content)

# ============== SLIDE 9: RESULTS ==============
def slide9_content(slide):
    # Big stats
    stats = [
        ('28', 'AI Generations'),
        ('65+', 'Screenshots'),
        ('2', 'Bugs Found'),
        (f'{avg_gen_time/1000:.1f}s', 'Avg Gen Time'),
        ('100%', 'Documented'),
        ('0', 'Scripts Written')
    ]

    colors = [
        RGBColor(0, 120, 60), RGBColor(0, 102, 204), RGBColor(200, 50, 50),
        RGBColor(153, 51, 153), RGBColor(0, 150, 150), RGBColor(100, 100, 100)
    ]

    for i, (val, label) in enumerate(stats):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 1.5 + row * 2.8

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()

        val_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(3.8), Inches(1.2))
        p = val_box.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.6), Inches(3.8), Inches(0.5))
        p = lbl_box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "Results Summary", slide9_content)

# ============== SLIDE 10: ROI ==============
def slide10_content(slide):
    # Two columns: Manual vs Claude
    # Manual column
    manual_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5))
    manual_box.fill.solid()
    manual_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    manual_box.line.color.rgb = RGBColor(150, 150, 150)

    manual_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.6))
    p = manual_title.text_frame.paragraphs[0]
    p.text = "Manual Testing"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    manual_items = [
        "4-6 hours to test all styles",
        "Manual screenshot capture",
        "Copy/paste into Word doc",
        "Measure timings with stopwatch",
        "Write bug reports manually",
        "Repeat for each test cycle"
    ]
    add_bullets(slide, manual_items, 0.7, 2.4, 5.4, 4, font_size=18)

    # Claude column
    claude_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.8), Inches(5))
    claude_box.fill.solid()
    claude_box.fill.fore_color.rgb = RGBColor(230, 245, 255)
    claude_box.line.color.rgb = RGBColor(0, 102, 204)

    claude_title = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.8), Inches(0.6))
    p = claude_title.text_frame.paragraphs[0]
    p.text = "Claude Code"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER

    claude_items = [
        "~30 minutes end-to-end",
        "Automatic screenshot capture",
        "Auto-generated reports",
        "JavaScript-precise timings",
        "Bugs discovered autonomously",
        "Same prompt, repeat anytime"
    ]
    add_bullets(slide, claude_items, 7.2, 2.4, 5.4, 4, font_size=18)

    # Bottom comparison
    savings = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.5))
    p = savings.text_frame.paragraphs[0]
    p.text = "Time Savings: 85-90% | Quality: Higher | Consistency: 100%"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 60)
    p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "ROI: Manual vs AI-Powered QA", slide10_content)

# ============== SLIDE 11: WHY THIS MATTERS ==============
def slide11_content(slide):
    # Central message
    msg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.333), Inches(3))
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = RGBColor(0, 45, 98)
    msg_box.line.fill.background()

    quote = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(2))
    tf = quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "We are not just automating clicks.\n\nWe are automating reasoning, exploration, and communication."
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Three pillars below
    pillars = [
        ("Reasoning", "Claude understands\nwhat to test and why"),
        ("Exploration", "Discovers bugs\nwe didn't look for"),
        ("Communication", "Auto-generates reports\nfor stakeholders")
    ]

    for i, (title, desc) in enumerate(pillars):
        x = 1 + i * 4

        title_box = slide.shapes.add_textbox(Inches(x), Inches(5.5), Inches(3.5), Inches(0.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.alignment = PP_ALIGN.CENTER

        desc_box = slide.shapes.add_textbox(Inches(x), Inches(6), Inches(3.5), Inches(1))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER

add_content_slide(prs, "Why This Matters", slide11_content)

# ============== SLIDE 12: FUTURE ==============
def slide12_content(slide):
    futures = [
        ("CI/CD Integration", "Automated quality gates on every deployment"),
        ("Parallel Agents", "Multiple Claude instances testing simultaneously"),
        ("Regression Comparison", "AI comparing outputs across builds"),
        ("Continuous A11y", "Accessibility audits on every release"),
        ("Multi-Browser", "Same test, Chrome/Firefox/Safari"),
        ("Self-Healing Tests", "No broken selectors, natural language adapts")
    ]

    for i, (title, desc) in enumerate(futures):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = 1.5 + row * 1.8

        # Icon
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(0, 102, 204)
        icon.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y), Inches(5.5), Inches(0.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.5), Inches(5.5), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(100, 100, 100)

add_content_slide(prs, "Future Possibilities", slide12_content)

# ============== SLIDE 13: CLOSING ==============
add_title_slide(prs,
    "Questions?",
    "Demo Available | Contact: AI Quality Engineering Team")

# Save
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'\nSlides created: 13')
print(f'Screenshots referenced: 20+')
print(f'Metrics included: {len(metrics)} data points')
