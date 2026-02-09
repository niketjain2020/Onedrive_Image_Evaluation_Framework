"""
Create Benchmark Story PowerPoint Presentation
Tells the story of the Restyle Benchmark Orchestrator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# Alias for convenience
RgbColor = RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, color=(0, 120, 215)):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Colored background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(*color)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, subtitle=None):
    """Add a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # Subtitle if provided
    y_offset = 1.2
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RgbColor(100, 100, 100)
        y_offset = 1.6

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # Left column title
    if left_title:
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 120, 215)

    # Right column title
    if right_title:
        right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 120, 215)

    y_start = 1.9 if left_title else 1.4

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_start), Inches(4.2), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(y_start), Inches(4.2), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.4), Inches(9), Inches(0.5 * num_rows)).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(0, 51, 102)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(255, 255, 255)

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(240, 240, 240)

    return slide

def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== TITLE =====
    add_title_slide(
        prs,
        "AI Restyle Benchmark",
        "Automated Quality Evaluation for OneDrive Photos AI Features\n" + datetime.now().strftime("%B %d, %Y")
    )

    # ===== THE PROBLEM =====
    add_section_slide(prs, "The Challenge", (0, 51, 102))

    add_content_slide(prs, "The Product Question", [
        "OneDrive Photos has 14 AI Restyle options",
        "Each transforms photos into different artistic styles",
        "But which styles should we feature prominently?",
        "Which ones actually delight users?",
        "How do we measure quality objectively?",
        "How do we detect regressions when models change?"
    ], subtitle="We needed data-driven answers, not opinions")

    add_two_column_slide(prs, "Manual Testing vs Automation",
        [
            "Subjective assessments",
            "Time-consuming (hours per style)",
            "Inconsistent across testers",
            "No historical tracking",
            "Hard to reproduce",
            "Limited coverage"
        ],
        [
            "Objective scoring framework",
            "Minutes per full benchmark",
            "Consistent evaluation criteria",
            "Full audit trail in Excel",
            "100% reproducible",
            "All styles, all images"
        ],
        "Before: Manual", "After: Automated"
    )

    # ===== THE SOLUTION =====
    add_section_slide(prs, "The Solution", (0, 120, 215))

    add_content_slide(prs, "Benchmark Orchestrator v1.0", [
        "Single source of truth: run_spec.json",
        "7-phase execution pipeline",
        "Two AI judges with different perspectives",
        "ACRUE v3 evaluation framework",
        "Automated artifact collection",
        "Excel ledger for historical tracking"
    ], subtitle="Fully automated, reproducible, data-driven")

    add_content_slide(prs, "The Competitors", [
        "Three styles enter the arena:",
        "",
        "ANIME - Big expressive eyes, vibrant colors",
        "POP ART - Bold Warhol vibes, halftone dots",
        "STORYBOOK - Soft watercolors, warm nostalgia",
        "",
        "Each applied to 3 diverse photos = 9 transformations"
    ])

    # ===== THE PROCESS =====
    add_section_slide(prs, "The 7-Phase Pipeline", (0, 153, 153))

    add_table_slide(prs, "Pipeline Overview",
        ["Phase", "Name", "What Happens"],
        [
            ["1", "VALIDATE", "Check all preconditions before starting"],
            ["2", "CAPTURE", "Select and screenshot 3 source images"],
            ["3", "TRANSFORM", "Apply each style to each image (9 total)"],
            ["4", "EVALUATE", "ACRUE v3 scoring + feasibility ranking"],
            ["5", "PREFERENCE", "Opus judges aesthetic appeal"],
            ["6", "SYNTHESIZE", "Combine rankings into final verdict"],
            ["7", "PERSIST", "Save artifacts, update Excel ledger"]
        ]
    )

    add_content_slide(prs, "Phase 1: Validate", [
        "8 precondition checks must ALL pass:",
        "",
        "V1-V2: run_spec.json exists and is valid JSON",
        "V3-V4: Exactly 3 styles and 3 images configured",
        "V5: Output directory doesn't already exist",
        "V6: Baseline run exists (if specified)",
        "V7: GEMINI_API_KEY is set",
        "V8: Playwright MCP is connected",
        "",
        "ANY failure = ABORT. No partial runs allowed."
    ])

    add_content_slide(prs, "Phase 2-3: Capture & Transform", [
        "Playwright MCP automates the browser:",
        "",
        "1. Navigate to OneDrive Photos",
        "2. Select 3 diverse images",
        "3. For each image, screenshot the original",
        "4. Click Edit > Restyle",
        "5. Apply each of 3 styles (~50s generation each)",
        "6. Screenshot each result",
        "",
        "Output: 3 originals + 9 restyled = 12 images"
    ])

    add_content_slide(prs, "Phase 4: ACRUE Evaluation", [
        "Gemini Flash evaluates each image on 5 dimensions:",
        "",
        "A - Accuracy: Does it preserve identity while applying style?",
        "C - Completeness: Is the style applied everywhere?",
        "R - Relevance: Does it actually look like the requested style?",
        "U - Usefulness: Is it free of artifacts and glitches?",
        "E - Exceptional: Does it have 'wow factor'?",
        "",
        "Each dimension has 4-5 Yes/No assertions + confidence scores"
    ])

    add_table_slide(prs, "ACRUE Scoring Weights",
        ["Dimension", "Weight", "Max Score", "Focus"],
        [
            ["Accuracy", "1.0", "5.0", "Identity preservation"],
            ["Completeness", "1.0", "5.0", "Full coverage"],
            ["Relevance", "0.5", "2.5", "Style authenticity"],
            ["Usefulness", "0.5", "2.5", "Practical quality"],
            ["Exceptional", "2.0", "10.0", "Delight factor"],
            ["TOTAL", "5.0", "25.0", ""]
        ]
    )

    add_content_slide(prs, "Phase 5: Opus Preference Judge", [
        "Claude Opus evaluates aesthetic appeal:",
        "",
        "Visual appeal / 'wow factor'",
        "User delight potential",
        "Social shareability",
        "'I want to frame this' quality",
        "",
        "Different perspective from technical ACRUE scores",
        "Captures the emotional/subjective dimension"
    ])

    add_content_slide(prs, "Phase 6-7: Synthesize & Persist", [
        "Combine both judges' rankings:",
        "",
        "final_score = 0.5 x gemini_rank + 0.5 x opus_rank",
        "Winner = lowest final score",
        "",
        "Generate report.md with full analysis",
        "Update Excel ledger with run summary",
        "Lock all artifacts for audit trail"
    ])

    # ===== ARTIFACTS =====
    add_section_slide(prs, "What We Produce", (102, 51, 153))

    add_content_slide(prs, "Artifact Manifest (16+ files)", [
        "runs/Run_2026_02_04_v1/",
        "  run_spec.json     - Frozen config",
        "  originals/        - 3 source images",
        "  restyled/         - 9 transformed images",
        "  acrue.json        - 9 detailed evaluations",
        "  gemini.json       - Feasibility rankings",
        "  opus.json         - Preference rankings",
        "  report.md         - Final synthesis",
        "",
        "Plus: ai_restyle_benchmark.xlsx updated"
    ])

    # ===== VALUE =====
    add_section_slide(prs, "The Value", (0, 153, 76))

    add_two_column_slide(prs, "What This Enables",
        [
            "Data-driven feature prioritization",
            "Objective style comparisons",
            "Regression detection",
            "Historical quality tracking",
            "Evidence for leadership",
            "Reproducible benchmarks"
        ],
        [
            "A/B testing AI features",
            "Prompt engineering eval",
            "Model comparison",
            "CI/CD quality gates",
            "Competitive analysis",
            "User study automation"
        ],
        "Current Use", "Future Applications"
    )

    # ===== FUTURE APPLICATIONS =====
    add_section_slide(prs, "Beyond Restyle", (153, 102, 0))

    add_table_slide(prs, "Other Use Cases",
        ["Use Case", "What You'd Compare", "Value"],
        [
            ["A/B Testing", "Old algo vs New algo", "Validate improvements"],
            ["Regression Testing", "Build N vs Build N+1", "Catch quality drops"],
            ["Prompt Engineering", "Prompt v1 vs v2 vs v3", "Optimize prompts"],
            ["Model Comparison", "DALL-E vs SD vs MJ", "Select best model"],
            ["Accessibility", "WCAG compliance checks", "Ensure inclusion"],
            ["CI/CD Gate", "Block deploys if <80%", "Automated QA"]
        ]
    )

    add_content_slide(prs, "The Pattern is Reusable", [
        "INPUT  ->  TRANSFORM  ->  JUDGE  ->  SYNTHESIZE  ->  RECORD",
        "",
        "Plug in ANY:",
        "  - Input source (OneDrive, local, URLs, APIs)",
        "  - Transformation (styles, filters, models)",
        "  - Judges (Gemini, Claude, GPT-4, humans)",
        "  - Synthesis formula (weighted, voting, veto)",
        "",
        "Framework adapts to many evaluation needs"
    ])

    # ===== SUMMARY =====
    add_section_slide(prs, "Summary", (0, 51, 102))

    add_content_slide(prs, "Key Takeaways", [
        "Automated quality evaluation for AI features",
        "7-phase pipeline with strict validation",
        "Two judges: technical (Gemini) + aesthetic (Opus)",
        "ACRUE v3 framework with 23 assertions per image",
        "Full artifact trail for reproducibility",
        "Extensible to other AI evaluation needs"
    ])

    add_content_slide(prs, "Current Status", [
        "Orchestrator v1.0 implemented and tested",
        "All validation checks passing",
        "Style assertions ready for Anime, Pop Art, Storybook",
        "Excel ledger initialized",
        "",
        "Ready to run first benchmark!",
        "",
        "Command: python benchmark_orchestrator.py"
    ])

    # Save
    output_path = "Restyle_Benchmark_Story.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
