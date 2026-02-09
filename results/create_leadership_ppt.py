from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

# Screenshot folder
screenshot_folder = str(Path(__file__).resolve().parent.parent.parent / '.playwright-mcp')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content area
    if image_path and os.path.exists(image_path):
        # With image - split layout
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
        slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.5), width=Inches(6.3))
    else:
        # Full width
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 60)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    return slide

def add_image_slide(prs, title, image_path, caption):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.2), width=Inches(10.333))

    # Caption
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    tf = cap_box.text_frame
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_quote_slide(prs, quote, attribution):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 240, 240)
    background.line.fill.background()

    # Quote mark
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(1), Inches(1))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(120)
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
    tf = attr_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {attribution}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.RIGHT

    return slide

# ============== CREATE SLIDES ==============

# Slide 1: Title
add_title_slide(prs,
    "Using AI to Test AI",
    "How Claude Code is Transforming Our QA Process for OneDrive Photos AI Restyle")

# Slide 2: The Challenge
add_content_slide(prs, "The Challenge: Testing AI-Powered Features at Scale", [
    "• OneDrive Photos AI Restyle: 14 style presets generating unique images",
    "• Each generation takes 60-90 seconds with non-deterministic outputs",
    "• Traditional automation tools struggle with:",
    "    - Dynamic UI elements that change during AI generation",
    "    - No predictable output to validate against",
    "    - Complex multi-step workflows (select → generate → undo → save)",
    "• Manual testing is time-consuming and inconsistent",
    "• Need: Intelligent agent that can reason, adapt, and validate behavior"
])

# Slide 3: Our Approach
add_two_column_slide(prs, "Our Approach: AI-Powered QA Agent",
    "Traditional Automation",
    [
        "Fixed scripts with hardcoded selectors",
        "Breaks when UI changes",
        "Cannot handle dynamic content",
        "No reasoning capability",
        "Requires constant maintenance",
        "Limited to predefined scenarios"
    ],
    "Claude Code + Playwright",
    [
        "Natural language test instructions",
        "Adapts to UI changes automatically",
        "Understands context and intent",
        "Reasons through complex scenarios",
        "Self-healing when issues occur",
        "Can explore and discover issues"
    ]
)

# Slide 4: Setup & Input
add_content_slide(prs, "Setup: What We Provided", [
    "• Claude Code CLI with Playwright MCP (Model Context Protocol)",
    "• Single natural language prompt describing the test scenario:",
    "",
    '  "Go to OneDrive Photos. Open first image. Click AI Restyle.',
    '   Select first style. Generate. Test Stop, Back, Undo, Redo,',
    '   Reset, Save, Copy, Download. Take screenshots. Generate report."',
    "",
    "• No selectors, no scripts, no test framework setup",
    "• Claude interprets intent and executes autonomously",
    "• Total setup time: < 5 minutes"
])

# Slide 5: Claude's Reasoning - Example 1
add_image_slide(prs, "Claude's Reasoning: Handling Authentication",
    os.path.join(screenshot_folder, 'qa_step0_gallery.png'),
    "Claude detected redirect to login, waited for auth, then continued to Photos gallery")

# Slide 6: Claude's Reasoning - Example 2
add_image_slide(prs, "Claude's Reasoning: Discovering UI Elements",
    os.path.join(screenshot_folder, 'qa_step2_restyle_panel.png'),
    "Claude identified all 14 style presets and documented them without being told what to expect")

# Slide 7: Claude's Reasoning - Example 3
add_content_slide(prs, "Claude's Reasoning: Adaptive Problem Solving", [
    "Challenge 1: Generation completed before Stop button could be tested",
    "  → Claude noted 'N/A' and documented Stop button was visible during generation",
    "",
    "Challenge 2: No explicit 'Copy' button in UI",
    "  → Claude reported as N/A, suggested it may be in context menu",
    "",
    "Challenge 3: After Reset, Save/Download buttons were disabled",
    "  → Claude understood it needed to regenerate to test these buttons",
    "  → Autonomously re-ran generation to complete remaining tests",
    "",
    "Challenge 4: Back button triggered unsaved changes dialog",
    "  → Claude captured dialog, documented UX behavior, then dismissed to continue"
])

# Slide 8: Generation Results
add_image_slide(prs, "AI Testing AI: Creative Output Verification",
    os.path.join(screenshot_folder, 'qa_step4_generation_complete.png'),
    '"STARFRUIT HUSTLE" - Claude verified generation success by checking UI state changes')

# Slide 9: State Validation
add_image_slide(prs, "Intelligent State Validation: Undo/Redo Testing",
    os.path.join(screenshot_folder, 'qa_step7a_after_undo.png'),
    "Claude verified Undo reverted to original AND checked button states changed correctly")

# Slide 10: Results
add_two_column_slide(prs, "Results: Comprehensive QA in One Session",
    "Test Execution",
    [
        "11 test scenarios executed",
        "10 passed, 0 failed, 2 N/A",
        "13 screenshots captured",
        "Full Word report generated",
        "Total time: ~5 minutes",
        "Zero script maintenance"
    ],
    "Findings Documented",
    [
        "All core flows working correctly",
        "UI state consistency verified",
        "Generation time benchmarked (~60s)",
        "UX behaviors captured (dialogs)",
        "Creative variation confirmed",
        "No visual glitches observed"
    ]
)

# Slide 11: Final Output
add_image_slide(prs, "Final Deliverable: Professional QA Report",
    os.path.join(screenshot_folder, 'qa_step9_saved.png'),
    "Claude generated a complete Word document with step-by-step results and embedded screenshots")

# Slide 12: Key Insight
add_quote_slide(prs,
    "We're using AI to test AI — Claude doesn't just execute tests, it understands what we're testing and why, adapting its approach when things don't go as expected.",
    "The Future of Quality Assurance")

# Slide 13: Possibilities
add_content_slide(prs, "Future Possibilities", [
    "• Parallel Testing: Multiple Claude agents testing different styles simultaneously",
    "• Regression Detection: AI comparing outputs across builds",
    "• Exploratory Testing: Claude discovering edge cases we haven't thought of",
    "• Cross-Browser Testing: Same natural language prompt, different browsers",
    "• Accessibility Testing: Claude evaluating WCAG compliance",
    "• Performance Monitoring: Tracking generation times across environments",
    "• Self-Documenting Tests: Every run produces human-readable reports",
    "• Integration with CI/CD: Automated quality gates with AI reasoning"
])

# Slide 14: Call to Action
add_content_slide(prs, "Getting Started: What You Need", [
    "• Claude Code CLI (Anthropic's official tool)",
    "• Playwright MCP for browser automation",
    "• Natural language descriptions of test scenarios",
    "",
    "Investment:",
    "  - Setup: 30 minutes (one-time)",
    "  - Per test: Natural language prompt + review results",
    "",
    "ROI:",
    "  - Reduced test maintenance burden",
    "  - Higher test coverage with less effort",
    "  - Documentation generated automatically",
    "  - Faster feedback on AI feature quality"
])

# Slide 15: Thank You
add_title_slide(prs,
    "Questions?",
    "Demo available | Contact: AI Quality Engineering Team")

# Save presentation
output_path = str(Path(__file__).resolve().parent / 'AI_Testing_AI_Leadership_Presentation.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
