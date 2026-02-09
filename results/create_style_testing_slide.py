from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Dark background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = RGBColor(15, 15, 25)
background.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(10), Inches(0.6))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "AI Restyle: Custom Style Prompt Discovery"
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.75), Inches(10), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.text = "Autonomous testing with Claude Code + Playwright MCP"
subtitle_para.font.size = Pt(16)
subtitle_para.font.color.rgb = RGBColor(150, 150, 180)

# Left section - What We Did
left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(3.8), Inches(2.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

p1 = left_frame.paragraphs[0]
p1.text = "What Claude Code Did:"
p1.font.size = Pt(14)
p1.font.bold = True
p1.font.color.rgb = RGBColor(100, 200, 255)

items = [
    "Navigated to OneDrive Photos",
    "Opened AI Restyle panel",
    "Tested 12 custom style prompts",
    "Captured 12 result screenshots",
    "Evaluated transformation quality",
    "Generated DOCX report with images",
    "Recommended top 3 styles"
]

for item in items:
    p = left_frame.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.space_before = Pt(4)

# Stats box
stats_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(4.2), Inches(3.8), Inches(1.1))
stats_box.fill.solid()
stats_box.fill.fore_color.rgb = RGBColor(40, 40, 60)
stats_box.line.color.rgb = RGBColor(100, 100, 140)

stats_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.35), Inches(3.4), Inches(0.9))
stats_frame = stats_text.text_frame
sp = stats_frame.paragraphs[0]
sp.text = "12 prompts tested  |  ~50 sec each"
sp.font.size = Pt(11)
sp.font.color.rgb = RGBColor(180, 180, 200)
sp.alignment = PP_ALIGN.CENTER

sp2 = stats_frame.add_paragraph()
sp2.text = "Fully autonomous - no manual intervention"
sp2.font.size = Pt(11)
sp2.font.bold = True
sp2.font.color.rgb = RGBColor(100, 255, 150)
sp2.alignment = PP_ALIGN.CENTER

# Right section - Top 3 Results with images
results_title = slide.shapes.add_textbox(Inches(4.3), Inches(1.4), Inches(8.5), Inches(0.4))
rt_frame = results_title.text_frame
rt_para = rt_frame.paragraphs[0]
rt_para.text = "TOP 3 RECOMMENDED STYLES"
rt_para.font.size = Pt(14)
rt_para.font.bold = True
rt_para.font.color.rgb = RGBColor(255, 200, 100)

screenshot_folder = str(Path(__file__).resolve().parent.parent.parent / '.playwright-mcp')

# Three result cards
results = [
    {
        'name': '#1 Japanese Ukiyo-e',
        'desc': 'Traditional woodblock art',
        'img': 'prompt05_japanese_ukiyoe.png'
    },
    {
        'name': '#2 Surrealist Dreamscape',
        'desc': 'Magical floating elements',
        'img': 'prompt08_surrealist_dreamscape.png'
    },
    {
        'name': '#3 Art Deco Poster',
        'desc': 'Geometric vintage design',
        'img': 'prompt09_art_deco_poster.png'
    }
]

x_positions = [4.3, 7.3, 10.3]

for i, result in enumerate(results):
    x = Inches(x_positions[i])

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(2.8), Inches(3.6))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 30, 45)
    card.line.color.rgb = RGBColor(80, 80, 120)

    # Add image
    img_path = os.path.join(screenshot_folder, result['img'])
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x + Inches(0.1), Inches(2.0), width=Inches(2.6))

    # Label
    label = slide.shapes.add_textbox(x, Inches(4.85), Inches(2.8), Inches(0.7))
    lf = label.text_frame
    lp = lf.paragraphs[0]
    lp.text = result['name']
    lp.font.size = Pt(12)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(255, 255, 255)
    lp.alignment = PP_ALIGN.CENTER

    lp2 = lf.add_paragraph()
    lp2.text = result['desc']
    lp2.font.size = Pt(10)
    lp2.font.color.rgb = RGBColor(180, 180, 200)
    lp2.alignment = PP_ALIGN.CENTER

# Bottom - Original image for reference
orig_label = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(3.8), Inches(0.3))
ol_frame = orig_label.text_frame
ol_para = ol_frame.paragraphs[0]
ol_para.text = "Original Test Image:"
ol_para.font.size = Pt(11)
ol_para.font.color.rgb = RGBColor(150, 150, 180)

# Add original image
orig_img = os.path.join(screenshot_folder, 'test1_adult_original.png')
if os.path.exists(orig_img):
    slide.shapes.add_picture(orig_img, Inches(0.3), Inches(5.85), width=Inches(2.2))

# Workflow arrow at bottom
workflow_box = slide.shapes.add_textbox(Inches(4.3), Inches(5.7), Inches(8.7), Inches(1.5))
wf_frame = workflow_box.text_frame
wf_frame.word_wrap = True

wp = wf_frame.paragraphs[0]
wp.text = "Workflow: Claude Code autonomously navigated browser, entered custom prompts,"
wp.font.size = Pt(10)
wp.font.color.rgb = RGBColor(150, 150, 180)

wp2 = wf_frame.add_paragraph()
wp2.text = "waited for AI generation (~50s each), captured screenshots, and compiled recommendations."
wp2.font.size = Pt(10)
wp2.font.color.rgb = RGBColor(150, 150, 180)

# Tools badge
tools_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(6.8), Inches(2.6), Inches(0.5))
tools_box.fill.solid()
tools_box.fill.fore_color.rgb = RGBColor(60, 60, 90)
tools_box.line.fill.background()

tools_text = slide.shapes.add_textbox(Inches(10.5), Inches(6.85), Inches(2.6), Inches(0.4))
tt_frame = tools_text.text_frame
tt_para = tt_frame.paragraphs[0]
tt_para.text = "Claude Code + Playwright"
tt_para.font.size = Pt(11)
tt_para.font.bold = True
tt_para.font.color.rgb = RGBColor(200, 200, 255)
tt_para.alignment = PP_ALIGN.CENTER

# Save
output_path = str(Path(__file__).resolve().parent / 'Style_Prompt_Discovery_Slide.pptx')
prs.save(output_path)
print(f'Slide saved to: {output_path}')
