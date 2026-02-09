from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

# Screenshot folder
screenshot_folder = str(Path(__file__).resolve().parent.parent.parent / '.playwright-mcp')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content area
    if image_path and os.path.exists(image_path):
        # With image - split layout
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
        slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.5), width=Inches(6.3))
    else:
        # Full width
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 60)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    return slide

def add_image_slide(prs, title, image_path, caption):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.2), width=Inches(10.333))

    # Caption
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    tf = cap_box.text_frame
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_stats_slide(prs, title, stats):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Stats boxes
    box_width = Inches(3.8)
    box_height = Inches(2.2)
    start_x = Inches(0.7)
    gap = Inches(0.5)

    colors = [RGBColor(0, 120, 60), RGBColor(0, 102, 204), RGBColor(153, 51, 153)]

    for i, (stat_value, stat_label) in enumerate(stats[:3]):
        x = start_x + i * (box_width + gap)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Stat value
        val_box = slide.shapes.add_textbox(x, Inches(2.0), box_width, Inches(1.2))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat_value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Stat label
        label_box = slide.shapes.add_textbox(x, Inches(3.2), box_width, Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat_label
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Second row if more than 3 stats
    if len(stats) > 3:
        for i, (stat_value, stat_label) in enumerate(stats[3:6]):
            x = start_x + i * (box_width + gap)

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.5), box_width, box_height)
            box.fill.solid()
            box.fill.fore_color.rgb = colors[(i + 3) % len(colors)]
            box.line.fill.background()

            val_box = slide.shapes.add_textbox(x, Inches(4.7), box_width, Inches(1.2))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = stat_value
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            label_box = slide.shapes.add_textbox(x, Inches(5.9), box_width, Inches(0.6))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = stat_label
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    return slide

def add_quote_slide(prs, quote, attribution):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 240, 240)
    background.line.fill.background()

    # Quote mark
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(1), Inches(1))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(120)
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
    tf = attr_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {attribution}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.RIGHT

    return slide

# ============== CREATE SLIDES ==============

# Slide 1: Title
add_title_slide(prs,
    "Claude Code + Playwright MCP",
    "Transforming QA, Documentation & Automation with AI-Powered Agents")

# Slide 2: What We Built
add_content_slide(prs, "What We Built: An AI-Powered QA Ecosystem", [
    "• Claude Code CLI - Anthropic's official AI coding assistant",
    "• Playwright MCP - Browser automation via Model Context Protocol",
    "• Natural Language Interface - No scripts, just describe what you want",
    "",
    "The Result:",
    "  → An autonomous agent that can test, audit, document, and report",
    "  → All from simple English instructions",
    "  → Adapts to UI changes without code maintenance"
])

# Slide 3: Use Cases Overview
add_stats_slide(prs, "12 Use Cases Demonstrated", [
    ("12", "Use Cases"),
    ("9", "Screenshots Auto-Captured"),
    ("4", "Reports Generated"),
    ("5", "A11y Bugs Found"),
    ("2", "Functional Bugs"),
    ("50.5s", "AI Gen Time Tracked")
])

# Slide 4: Use Case 1 - Functional Testing
add_image_slide(prs, "Use Case 1: Functional QA Testing",
    os.path.join(screenshot_folder, 'audit_03_restyle_panel.png'),
    "Claude autonomously navigated to OneDrive Photos, opened images, and tested AI Restyle feature end-to-end")

# Slide 5: Functional Testing Details
add_content_slide(prs, "Functional Testing: What Claude Tested", [
    "• Navigation to OneDrive Photos gallery",
    "• Opening images in viewer",
    "• AI Restyle feature - all 14 style presets verified",
    "• Style selection and prompt population",
    "• AI generation workflow (start → progress → complete)",
    "• Undo/Redo/Reset functionality",
    "• Save/Download button states",
    "• Unsaved changes dialog flow",
    "• Format support (PNG, JPEG, WEBP)"
])

# Slide 6: Use Case 2 - Bug Discovery
add_two_column_slide(prs, "Use Case 2: Autonomous Bug Discovery",
    "Bug #1: WEBP Not Supported",
    [
        "WEBP images missing Restyle button",
        "No user messaging explaining why",
        "Severity: Medium (P2)",
        "Found without being told to look"
    ],
    "Bug #2: Stop Button Broken",
    [
        "Stop button visible during generation",
        "Clicking does NOT cancel operation",
        "Clicked 3 times - generation continued",
        "Severity: High (P1)"
    ]
)

# Slide 7: Use Case 3 - Latency Tracking
add_content_slide(prs, "Use Case 3: Performance & Latency Tracking", [
    "• JavaScript timestamps captured at each action",
    "• Precise measurement of user-perceived latency",
    "",
    "Results from our test session:",
    "  → Gallery Page Load: 911ms (Good)",
    "  → Restyle Panel Open: ~500ms (Good)",
    "  → AI Generation: 50.5 seconds (Acceptable for AI)",
    "  → Undo/Redo: <100ms (Excellent)",
    "",
    "• No manual stopwatch needed - automated tracking"
])

# Slide 8: Use Case 4 - Accessibility Auditing
add_image_slide(prs, "Use Case 4: Accessibility Auditing (WCAG)",
    os.path.join(screenshot_folder, 'audit_04_keyboard_nav_issue.png'),
    "Found 5 accessibility issues including keyboard navigation failure on style presets")

# Slide 9: Accessibility Findings
add_two_column_slide(prs, "Accessibility Audit: Findings",
    "Issues Found (5)",
    [
        "48 gallery images missing alt text",
        "Style presets not keyboard navigable",
        "Missing main landmark in gallery",
        "Restyle dialog missing aria-label",
        "Some buttons lack aria-labels"
    ],
    "Passes Verified (7)",
    [
        "Skip to main content link present",
        "Toolbar buttons have aria-labels",
        "Style preset images have alt text",
        "Unsaved changes dialog accessible",
        "Focus indicators visible",
        "Prompt textbox labeled",
        "Color contrast sufficient"
    ]
)

# Slide 10: Use Case 5 - Screenshot Capture
add_content_slide(prs, "Use Case 5: Automated Screenshot Capture", [
    "• 9 screenshots captured automatically during testing",
    "• Each screenshot documents a specific test state:",
    "",
    "  audit_01 - Keyboard focus indicator on gallery",
    "  audit_02 - Image viewer with toolbar",
    "  audit_03 - Restyle panel with presets",
    "  audit_04 - Keyboard navigation issue evidence",
    "  audit_05 - Generation in progress state",
    "  audit_06 - Generation complete (Pop Art result)",
    "  audit_07 - After Undo operation",
    "  audit_08 - After Redo operation",
    "  audit_09 - Unsaved changes dialog"
])

# Slide 11: Use Case 6 - Report Generation
add_image_slide(prs, "Use Case 6: Automated Report Generation",
    os.path.join(screenshot_folder, 'audit_06_generation_complete.png'),
    "Claude generates professional DOCX reports with embedded screenshots, findings, and recommendations")

# Slide 12: Reports Generated
add_content_slide(prs, "Reports Generated Automatically", [
    "• AI_Restyle_Bug_Bash_Report.md - Quick markdown summary",
    "",
    "• AI_Restyle_Bug_Bash_Report.docx - Formal bug report with:",
    "    - Executive summary with metrics",
    "    - Detailed bug descriptions with steps to reproduce",
    "    - Embedded screenshots as evidence",
    "    - Recommendations prioritized by severity",
    "",
    "• Comprehensive_Bug_Bash_Report.docx - Full audit with:",
    "    - Latency measurements",
    "    - Accessibility findings (WCAG references)",
    "    - UI/UX critique",
    "    - Design improvement suggestions",
    "",
    "• PowerPoint presentations for leadership communication"
])

# Slide 13: Use Case 7 - UI/UX Evaluation
add_content_slide(prs, "Use Case 7: UI/UX Evaluation & Critique", [
    "Claude analyzed the user experience and provided feedback:",
    "",
    "Issues Identified:",
    "  → No progress indicator during 50s generation wait",
    "  → Stop button present but non-functional (misleading)",
    "  → WEBP format unsupported with no user messaging",
    "",
    "Positive UX Elements Noted:",
    "  → Clear unsaved changes dialog messaging",
    "  → Proper Undo/Redo button state management",
    "  → Visual style previews help set expectations",
    "  → AI disclaimer sets appropriate expectations"
])

# Slide 14: Use Case 8 - Design Suggestions
add_content_slide(prs, "Use Case 8: Design Improvement Suggestions", [
    "Claude provided 6 actionable design recommendations:",
    "",
    "1. Add progress bar during AI generation",
    "2. Implement keyboard navigation for style presets",
    "   (roving tabindex pattern suggested)",
    "3. Show estimated generation time before starting",
    "4. Add format support indicator for unsupported images",
    "5. Add keyboard shortcuts (Ctrl+Z, Ctrl+S, Escape)",
    "6. Implement AI-generated alt text for gallery images",
    "",
    "Each suggestion includes implementation guidance"
])

# Slide 15: Use Case 9-12
add_two_column_slide(prs, "Additional Use Cases",
    "Use Case 9: Video Recording",
    [
        "Browser session recording",
        "Full user journey capture",
        "Playback for stakeholder review",
        "Evidence for complex bugs"
    ],
    "Use Cases 10-12",
    [
        "Email-ready report formatting",
        "Exploratory edge case testing",
        "Cross-format compatibility testing",
        "Regression validation"
    ]
)

# Slide 16: The Power of Natural Language
add_content_slide(prs, "The Power: Natural Language Instructions", [
    "What we typed:",
    "",
    '  "Go to OneDrive Photos. Test the AI Restyle feature.',
    '   Track latency at each action. Audit accessibility.',
    '   Critique the UI/UX. Suggest design improvements.',
    '   Take screenshots. Generate a comprehensive report."',
    "",
    "What Claude did:",
    "  → Navigated, tested, measured, audited, critiqued",
    "  → Captured 9 screenshots as evidence",
    "  → Found 5 accessibility bugs + 2 functional bugs",
    "  → Generated 4 professional reports",
    "  → All in ~10 minutes with zero script maintenance"
])

# Slide 17: ROI Summary
add_stats_slide(prs, "Return on Investment", [
    ("< 5 min", "Setup Time"),
    ("~10 min", "Full Test Suite"),
    ("0", "Scripts to Maintain"),
    ("7", "Bugs Found"),
    ("4", "Reports Generated"),
    ("100%", "Documented")
])

# Slide 18: Key Insight Quote
add_quote_slide(prs,
    "We didn't write test scripts. We described what we wanted to test in plain English, and Claude became our QA engineer, accessibility auditor, UX reviewer, and technical writer — all in one session.",
    "The Future of QA Automation")

# Slide 19: What's Next
add_content_slide(prs, "Future Possibilities", [
    "• Parallel Testing - Multiple Claude agents testing simultaneously",
    "• CI/CD Integration - Automated quality gates with AI reasoning",
    "• Regression Detection - AI comparing outputs across builds",
    "• Multi-Browser Testing - Same prompt, different browsers",
    "• Continuous Accessibility - A11y audits on every deployment",
    "• Self-Healing Tests - No more broken selectors",
    "• Intelligent Exploration - AI finding bugs we haven't thought of",
    "• Auto-Documentation - Every test run produces stakeholder reports"
])

# Slide 20: Call to Action
add_content_slide(prs, "Getting Started", [
    "Requirements:",
    "  • Claude Code CLI (Anthropic)",
    "  • Playwright MCP for browser automation",
    "  • Natural language test descriptions",
    "",
    "Investment:",
    "  • One-time setup: ~30 minutes",
    "  • Per test: Describe in English → Review results",
    "",
    "Benefits:",
    "  • Eliminate test script maintenance",
    "  • Higher coverage with less effort",
    "  • Auto-generated documentation",
    "  • Accessible to non-developers"
])

# Slide 21: Thank You
add_title_slide(prs,
    "Questions?",
    "Demo Available | Contact: AI Quality Engineering Team")

# Save presentation
output_path = str(Path(__file__).resolve().parent / 'Claude_Code_Use_Cases_Presentation.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
