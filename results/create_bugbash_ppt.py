from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

# Screenshot folder
screenshot_folder = str(Path(__file__).resolve().parent.parent.parent / '.playwright-mcp')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 25, 112)  # Midnight blue
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(25, 25, 112)
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content area
    if image_path and os.path.exists(image_path):
        # With image - split layout
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
        slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.5), width=Inches(6.3))
    else:
        # Full width
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(25, 25, 112)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 20, 60)  # Crimson for bugs

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)  # Green for passed

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

    return slide

def add_image_slide(prs, title, image_path, caption):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(25, 25, 112)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.2), width=Inches(10.333))

    # Caption
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    tf = cap_box.text_frame
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_quote_slide(prs, quote, attribution):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 240, 240)
    background.line.fill.background()

    # Quote mark
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(1), Inches(1))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(120)
    p.font.color.rgb = RGBColor(25, 25, 112)

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
    tf = attr_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {attribution}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 25, 112)
    p.alignment = PP_ALIGN.RIGHT

    return slide

def add_stats_slide(prs, title, stats):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(25, 25, 112)
    title_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Stats boxes
    box_width = Inches(2.8)
    box_height = Inches(2)
    start_x = Inches(0.8)
    y_pos = Inches(2.5)
    gap = Inches(0.4)

    colors = [
        RGBColor(65, 105, 225),   # Royal Blue
        RGBColor(0, 128, 0),       # Green
        RGBColor(220, 20, 60),     # Crimson
        RGBColor(255, 140, 0)      # Orange
    ]

    for i, (label, value) in enumerate(stats):
        x_pos = start_x + (box_width + gap) * i

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Value
        val_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.3), box_width, Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(1.3), box_width, Inches(0.6))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide

# ============== CREATE SLIDES ==============

# Slide 1: Title
add_title_slide(prs,
    "Autonomous Bug Hunting with AI",
    "How Claude Code Found 2 Critical Bugs in 10 Minutes — Without Human Intervention")

# Slide 2: The Challenge
add_content_slide(prs, "The Challenge: Bug Bash at Scale", [
    "• Traditional bug bashes require:",
    "    - Manual test case creation and documentation",
    "    - Human testers executing scenarios one by one",
    "    - Note-taking and screenshot capture by hand",
    "    - Hours of focused testing time",
    "",
    "• What if we could:",
    "    - Give an AI agent a simple instruction: 'Bug bash this feature'",
    "    - Let it autonomously explore, test, and document",
    "    - Find bugs humans might miss",
    "    - Generate professional reports automatically"
])

# Slide 3: The Experiment
add_content_slide(prs, "The Experiment: One Prompt, Zero Guidance", [
    "• Input given to Claude Code:",
    "",
    '  "Bug bash the AI Restyle feature"',
    "",
    "• What was NOT provided:",
    "    - No test cases or scenarios",
    "    - No list of formats to test",
    "    - No edge cases to explore",
    "    - No expected behaviors documented",
    "",
    "• Claude was left to figure out:",
    "    - What to test and how to test it",
    "    - What constitutes a bug vs. expected behavior",
    "    - How to document findings professionally"
])

# Slide 4: Claude's Autonomous Approach
add_content_slide(prs, "Claude's Autonomous Testing Strategy", [
    "• Claude self-organized into 5 test categories:",
    "",
    "  1. Input Edge Cases — Format support, file types",
    "  2. Generation Stress Tests — Double-clicks, rapid actions",
    "  3. State Management — Undo/Redo, button states",
    "  4. UI/UX Issues — Dialogs, confirmations, flow",
    "  5. Style Preset Verification — All 14 styles",
    "",
    "• No human told Claude to organize this way",
    "• Claude reasoned about comprehensive coverage autonomously"
])

# Slide 5: Test 1 - Format Discovery
add_image_slide(prs, "Bug Discovery #1: WEBP Format Testing",
    os.path.join(screenshot_folder, 'bugbash_01_webp_no_restyle.png'),
    "Claude autonomously tested different file formats and discovered WEBP is not supported")

# Slide 6: Bug #1 Details
add_content_slide(prs, "Bug #1: WEBP Format Not Supported", [
    "• What Claude did:",
    "    - Noticed gallery had multiple file formats (PNG, JPEG, WEBP)",
    "    - Decided to test each format's compatibility",
    "    - Opened a WEBP image and inspected the toolbar",
    "",
    "• What Claude found:",
    "    - WEBP images missing 'Edit' and 'Restyle with AI' buttons",
    "    - Only basic options available (Close, Delete, Download)",
    "",
    "• Why this matters:",
    "    - WEBP is a common modern format",
    "    - Users won't understand why feature is unavailable",
    "    - No error message explains the limitation"
])

# Slide 7: Test 2 - Stop Button
add_image_slide(prs, "Bug Discovery #2: Stop Button Testing",
    os.path.join(screenshot_folder, 'bugbash_02_stop_not_working.png'),
    "Claude clicked Stop 3 times during generation — it never stopped")

# Slide 8: Bug #2 Details
add_content_slide(prs, "Bug #2: Stop Button Doesn't Work", [
    "• What Claude did:",
    "    - Started an AI generation (Pop Art style)",
    "    - Clicked the Stop button during 'Pixels getting warmed up...'",
    "    - Clicked Stop again... and again (3 total clicks)",
    "",
    "• What Claude observed:",
    "    - Generation continued through all loading phases",
    "    - 'Brewing something cool...' → 'Pulling in those final bits...'",
    "    - Image was successfully generated despite Stop attempts",
    "",
    "• Claude's assessment:",
    "    - Severity: HIGH (P1)",
    "    - Stop button creates false expectation of control",
    "    - Users cannot cancel accidental generations"
])

# Slide 9: What Else Claude Tested
add_two_column_slide(prs, "Comprehensive Testing: What Claude Tried",
    "Stress Tests",
    [
        "Rapid style switching (3 styles in seconds)",
        "Double-click on Send button",
        "Custom prompt entry",
        "Back button with unsaved changes",
        "Discard functionality",
        "Reset after generation"
    ],
    "Verification Tests",
    [
        "All 14 style presets visible",
        "Button state management",
        "Undo/Redo functionality",
        "Save copy workflow",
        "Unsaved changes dialog",
        "Format support (PNG, JPEG, WEBP)"
    ]
)

# Slide 10: Unsaved Changes Test
add_image_slide(prs, "UX Validation: Unsaved Changes Dialog",
    os.path.join(screenshot_folder, 'bugbash_03_unsaved_dialog.png'),
    "Claude verified that the 'Leave without saving?' dialog works correctly — Good UX!")

# Slide 11: Results Summary
add_stats_slide(prs, "Bug Bash Results: By The Numbers", [
    ("Total Tests", "12"),
    ("Tests Passed", "10"),
    ("Bugs Found", "2"),
    ("Time Taken", "10 min")
])

# Slide 12: Time Savings
add_two_column_slide(prs, "Time Savings: Agentic vs. Manual",
    "Manual Bug Bash",
    [
        "Test case creation: 30-60 min",
        "Executing 12 scenarios: 45-60 min",
        "Screenshot capture: 15-20 min",
        "Note taking: 20-30 min",
        "Report writing: 30-45 min",
        "Total: 2.5-3.5 hours",
        "",
        "Requires: Dedicated tester focus"
    ],
    "Agentic Bug Bash",
    [
        "Prompt input: 30 seconds",
        "Autonomous execution: 10 min",
        "Screenshots: Automatic",
        "Documentation: Automatic",
        "Report generation: 1 min",
        "Total: ~12 minutes",
        "",
        "Requires: One sentence instruction"
    ]
)

# Slide 13: Key Insight - Autonomous Reasoning
add_content_slide(prs, "Key Insight: Autonomous Reasoning in Action", [
    "• Claude wasn't told to test WEBP files",
    "    → It noticed the format in the gallery and decided to test it",
    "",
    "• Claude wasn't told to click Stop multiple times",
    "    → It recognized the first click didn't work and persisted",
    "",
    "• Claude wasn't given severity ratings",
    "    → It reasoned that Stop button bug was HIGH priority",
    "",
    "• Claude organized its own test categories",
    "    → Created structure without templates or guidance",
    "",
    "• This is true agentic behavior — not scripted automation"
])

# Slide 14: What This Means
add_quote_slide(prs,
    "We gave Claude one instruction and walked away. It came back with 2 bugs, 12 documented tests, 4 screenshots, and a professional report. This is the future of quality assurance.",
    "The Power of Agentic AI")

# Slide 15: Deliverables Generated
add_content_slide(prs, "Auto-Generated Deliverables", [
    "• Without any templates or guidance, Claude produced:",
    "",
    "  1. Markdown Bug Report — Full technical documentation",
    "  2. Word Document (.docx) — Professional report with images",
    "  3. Screenshots — 4 key moments captured automatically",
    "  4. Severity Ratings — P1 and P2 classifications",
    "  5. Reproduction Steps — Detailed step-by-step instructions",
    "  6. Recommendations — Prioritized fix suggestions",
    "",
    "• All in < 15 minutes total, including report generation"
])

# Slide 16: Scaling Possibilities
add_content_slide(prs, "Scaling Autonomous Bug Hunting", [
    "• Parallel Testing: Multiple Claude agents testing different features",
    "• Nightly Bug Bashes: Automated quality sweeps on every build",
    "• Regression Detection: AI comparing behavior across versions",
    "• Exploratory Testing: Discovering edge cases humans don't think of",
    "• Cross-Platform: Same prompt, different environments",
    "• Continuous Documentation: Every test run produces artifacts",
    "",
    "• Investment: Natural language prompts",
    "• Return: Comprehensive quality assurance at scale"
])

# Slide 17: Conclusion
add_content_slide(prs, "Conclusion: The Agentic Advantage", [
    "• What we learned:",
    "    - AI can autonomously plan and execute bug bashes",
    "    - It finds bugs humans might overlook (format edge cases)",
    "    - It documents findings professionally without templates",
    "    - It saves 90%+ of traditional bug bash time",
    "",
    "• The shift:",
    "    - From: 'Write test cases, execute manually, document findings'",
    "    - To: 'Bug bash this feature' — and review the results",
    "",
    "• This is AI augmenting human testers, not replacing them",
    "    - Humans focus on judgment and priorities",
    "    - AI handles execution and documentation"
])

# Slide 18: Thank You
add_title_slide(prs,
    "Questions?",
    "Demo available | Full report: AI_Restyle_Bug_Bash_Report.docx")

# Save presentation
output_path = str(Path(__file__).resolve().parent / 'Autonomous_Bug_Bash_Presentation.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
