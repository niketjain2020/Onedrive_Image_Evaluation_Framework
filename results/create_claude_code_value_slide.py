from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Dark gradient-style background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = RGBColor(12, 12, 20)
background.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "Claude Code + Playwright: Autonomous QA Testing"
title_para.font.size = Pt(34)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

# Subtitle - value prop
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.text = "From manual bug bash to automated discovery, documentation, and reporting"
subtitle_para.font.size = Pt(16)
subtitle_para.font.color.rgb = RGBColor(150, 200, 255)

# ============ FOUR CAPABILITY CARDS ============

capabilities = [
    {
        'title': 'Functional Testing',
        'icon': 'BUG',
        'color': RGBColor(255, 100, 100),
        'items': [
            'End-to-end UI flows',
            'Button/interaction testing',
            'Error state validation',
            'Found: Stop button bug'
        ],
        'highlight': '2 bugs found autonomously'
    },
    {
        'title': 'Accessibility Audit',
        'icon': 'A11Y',
        'color': RGBColor(100, 200, 255),
        'items': [
            'WCAG compliance checks',
            'Missing alt text detection',
            'Keyboard navigation',
            'ARIA label validation'
        ],
        'highlight': '3 a11y issues identified'
    },
    {
        'title': 'Performance Tracking',
        'icon': 'PERF',
        'color': RGBColor(100, 255, 150),
        'items': [
            'Page load timing',
            'AI generation latency',
            'Action response times',
            'Baseline metrics captured'
        ],
        'highlight': '28 timing data points'
    },
    {
        'title': 'Auto Documentation',
        'icon': 'DOC',
        'color': RGBColor(255, 200, 100),
        'items': [
            'Screenshot capture',
            'Word reports with images',
            'PowerPoint generation',
            'Email-ready formats'
        ],
        'highlight': '65+ screenshots captured'
    }
]

card_width = 2.9
card_height = 2.6
start_x = 0.5
gap = 0.25

for i, cap in enumerate(capabilities):
    x = Inches(start_x + i * (card_width + gap))
    y = Inches(1.5)

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(card_width), Inches(card_height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(25, 25, 40)
    card.line.color.rgb = cap['color']
    card.line.width = Pt(2)

    # Title with colored accent
    title_tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(card_width - 0.3), Inches(0.4))
    tf = title_tb.text_frame
    tp = tf.paragraphs[0]
    tp.text = cap['title']
    tp.font.size = Pt(14)
    tp.font.bold = True
    tp.font.color.rgb = cap['color']

    # Bullet items
    items_tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), Inches(card_width - 0.3), Inches(1.4))
    itf = items_tb.text_frame
    itf.word_wrap = True

    for j, item in enumerate(cap['items']):
        if j == 0:
            p = itf.paragraphs[0]
        else:
            p = itf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(200, 200, 210)
        p.space_before = Pt(2)

    # Highlight stat box
    stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(2.05), Inches(card_width - 0.2), Inches(0.4))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = RGBColor(40, 40, 60)
    stat_box.line.fill.background()

    stat_tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(2.1), Inches(card_width - 0.2), Inches(0.35))
    stf = stat_tb.text_frame
    stp = stf.paragraphs[0]
    stp.text = cap['highlight']
    stp.font.size = Pt(10)
    stp.font.bold = True
    stp.font.color.rgb = cap['color']
    stp.alignment = PP_ALIGN.CENTER

# ============ BOTTOM SECTION - BEFORE/AFTER ============

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(60, 60, 80)
line.line.fill.background()

# Before/After comparison
before_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(5.5), Inches(0.4))
bt_frame = before_title.text_frame
bt_para = bt_frame.paragraphs[0]
bt_para.text = "BEFORE: Manual Bug Bash"
bt_para.font.size = Pt(14)
bt_para.font.bold = True
bt_para.font.color.rgb = RGBColor(255, 100, 100)

before_items = [
    "Hours of manual clicking through UI",
    "Screenshots saved manually to folders",
    "Copy/paste findings into Word docs",
    "Separate accessibility audit tools",
    "No timing data unless instrumented"
]

before_tb = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(5.5), Inches(1.8))
bf = before_tb.text_frame
bf.word_wrap = True
for j, item in enumerate(before_items):
    if j == 0:
        p = bf.paragraphs[0]
    else:
        p = bf.add_paragraph()
    p.text = f"✗  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(180, 180, 190)
    p.space_before = Pt(3)

# After
after_title = slide.shapes.add_textbox(Inches(6.8), Inches(4.5), Inches(6), Inches(0.4))
at_frame = after_title.text_frame
at_para = at_frame.paragraphs[0]
at_para.text = "AFTER: Claude Code + Playwright"
at_para.font.size = Pt(14)
at_para.font.bold = True
at_para.font.color.rgb = RGBColor(100, 255, 150)

after_items = [
    "Autonomous navigation & interaction",
    "Auto screenshot at every step",
    "DOCX/PPTX reports generated instantly",
    "Built-in accessibility snapshot checks",
    "JavaScript timing instrumentation"
]

after_tb = slide.shapes.add_textbox(Inches(6.8), Inches(4.9), Inches(6), Inches(1.8))
af = after_tb.text_frame
af.word_wrap = True
for j, item in enumerate(after_items):
    if j == 0:
        p = af.paragraphs[0]
    else:
        p = af.add_paragraph()
    p.text = f"✓  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 255, 200)
    p.space_before = Pt(3)

# Arrow between
arrow_tb = slide.shapes.add_textbox(Inches(6.0), Inches(5.4), Inches(0.7), Inches(0.5))
arrow_f = arrow_tb.text_frame
arrow_p = arrow_f.paragraphs[0]
arrow_p.text = "→"
arrow_p.font.size = Pt(36)
arrow_p.font.color.rgb = RGBColor(100, 200, 255)
arrow_p.alignment = PP_ALIGN.CENTER

# ============ BOTTOM STATS BAR ============

stats_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6))
stats_bar.fill.solid()
stats_bar.fill.fore_color.rgb = RGBColor(30, 30, 50)
stats_bar.line.fill.background()

stats_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.5))
st_frame = stats_text.text_frame
st_para = st_frame.paragraphs[0]
st_para.text = "OneDrive Photos AI Restyle  •  65+ screenshots  •  2 functional bugs  •  3 a11y issues  •  28 perf metrics  •  5 reports generated"
st_para.font.size = Pt(12)
st_para.font.color.rgb = RGBColor(180, 180, 200)
st_para.alignment = PP_ALIGN.CENTER

# Save
output_path = str(Path(__file__).resolve().parent / 'Claude_Code_QA_Value_Slide.pptx')
prs.save(output_path)
print(f'Slide saved to: {output_path}')
