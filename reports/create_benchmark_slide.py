"""
Generate a single slide: Why We Built an Automated AI Restyle Benchmark
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Alias for consistency
RgbColor = RGBColor

def add_section_title(slide, text, left, top, width, font_size=14, bold=True):
    """Add a section title"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RgbColor(0, 51, 102)  # Dark blue
    return shape

def add_bullet_list(slide, items, left, top, width, height, font_size=11):
    """Add a bullet list"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RgbColor(51, 51, 51)
        p.space_after = Pt(4)
    return shape

def add_flow_box(slide, text, left, top, width=1.3, height=0.5):
    """Add a flow diagram box"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0, 120, 212)  # Microsoft blue
    shape.line.color.rgb = RgbColor(0, 90, 158)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_arrow(slide, left, top):
    """Add a right arrow"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(0.3), Inches(0.25)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(100, 100, 100)
    shape.line.fill.background()
    return shape

def create_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(250, 250, 252)
    background.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(0, 51, 102)
    title_bar.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Claude Code + Eval Skill + Playwright + Gemini + Opus 4.5"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.5), Inches(12), Inches(0.3))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Why We Built an Automated AI Restyle Benchmark"
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(200, 220, 255)

    # ===== LEFT COLUMN: Problems =====
    add_section_title(slide, "The Problem with Human Evaluation", 0.4, 1.1, 4.5)

    problems = [
        "Evaluations are slow and don't scale as styles/prompts increase",
        "Results change based on who is in the room",
        "No durable benchmark to compare past vs current",
        "Manual process makes it hard to reuse results or track progress",
        "No systematic way to use multiple LLMs as judges"
    ]
    add_bullet_list(slide, problems, 0.4, 1.5, 4.3, 2.2, font_size=11)

    # ===== MIDDLE COLUMN: What Changes =====
    add_section_title(slide, "What This System Changes", 4.9, 1.1, 4.5)

    changes = [
        "Replaces ad-hoc human evals with standardized, repeatable pipeline",
        "Creates persistent benchmark that accumulates over time",
        "Plug in and compare different LLM judges as models evolve",
        "Separates feasibility from creative preference",
        'Enables statements like: "This prompt is 20% better"'
    ]
    add_bullet_list(slide, changes, 4.9, 1.5, 4.3, 2.2, font_size=11)

    # ===== RIGHT COLUMN: Outcome =====
    add_section_title(slide, "Outcome", 9.4, 1.1, 3.5)

    outcome_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.4), Inches(1.5), Inches(3.5), Inches(2.2)
    )
    outcome_box.fill.solid()
    outcome_box.fill.fore_color.rgb = RgbColor(232, 245, 233)  # Light green
    outcome_box.line.color.rgb = RgbColor(76, 175, 80)

    outcome_text = slide.shapes.add_textbox(Inches(9.5), Inches(1.6), Inches(3.3), Inches(2.0))
    tf = outcome_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Faster iteration, consistent decisions, and a defensible way to measure improvement as our AI restyle system evolves."
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(33, 37, 41)
    p.font.bold = True

    # ===== BOTTOM: High-Level Flow =====
    add_section_title(slide, "High-Level Flow", 0.4, 4.0, 12, font_size=14)

    # Flow diagram
    flow_y = 4.5

    # Box 1: Discover
    add_flow_box(slide, "Auto-discover\ntrending styles", 0.5, flow_y, 1.8, 0.7)
    add_arrow(slide, 2.4, flow_y + 0.22)

    # Box 2: Apply
    add_flow_box(slide, "Apply to diverse\nreal user photos", 2.8, flow_y, 1.8, 0.7)
    add_arrow(slide, 4.7, flow_y + 0.22)

    # Box 3: Evaluate
    add_flow_box(slide, "Evaluate with\nassertions + scoring", 5.1, flow_y, 1.8, 0.7)
    add_arrow(slide, 7.0, flow_y + 0.22)

    # Box 4: Rank
    add_flow_box(slide, "Rank with dual\nLLM judges", 7.4, flow_y, 1.8, 0.7)
    add_arrow(slide, 9.3, flow_y + 0.22)

    # Box 5: Persist
    add_flow_box(slide, "Persist to Excel\nfor comparison", 9.7, flow_y, 1.8, 0.7)

    # Labels under boxes
    labels = [
        ("Playwright scrapes\ntrend data", 0.5),
        ("Automated\nrestyle pipeline", 2.8),
        ("ACRUE v3\nframework", 5.1),
        ("Gemini + Opus\ncross-validation", 7.4),
        ("Longitudinal\ntracking", 9.7)
    ]

    for label_text, x in labels:
        label = slide.shapes.add_textbox(Inches(x), Inches(flow_y + 0.8), Inches(1.8), Inches(0.5))
        tf = label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label_text
        p.font.size = Pt(9)
        p.font.color.rgb = RgbColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    # Footer with tech stack icons (text representation)
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(12), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Tech Stack:  Claude Code  •  Playwright MCP  •  Gemini 2.0 Flash  •  Opus 4.5  •  ACRUE v3 Eval Framework  •  python-pptx"
    p.font.size = Pt(10)
    p.font.color.rgb = RgbColor(128, 128, 128)

    # Save
    output_path = "Automated_Restyle_Benchmark_Slide.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_slide()
