"""
Create Benchmark Story PowerPoint - Simple/End-User Version
Plain language that anyone can understand
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

RgbColor = RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title, color=(0, 120, 215)):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(*color)
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, subtitle=None):
    """Add a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    y_offset = 1.2
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RgbColor(100, 100, 100)
        y_offset = 1.6

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.space_after = Pt(14)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    if left_title:
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 120, 215)

    if right_title:
        right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 120, 215)

    y_start = 1.9 if left_title else 1.4

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_start), Inches(4.2), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(y_start), Inches(4.2), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    return slide

def add_big_quote_slide(prs, quote, attribution=""):
    """Add a big quote slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = RgbColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    if attribution:
        attr_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
        tf = attr_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"— {attribution}"
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_emoji_content_slide(prs, title, items_with_emoji):
    """Add a content slide with emoji bullets for visual appeal."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items_with_emoji):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.space_after = Pt(16)

    return slide

def create_presentation():
    """Create the simplified presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== TITLE =====
    add_title_slide(
        prs,
        "How Do We Know Which\nAI Style is Best?",
        "A smart way to test OneDrive's photo magic\n" + datetime.now().strftime("%B %d, %Y")
    )

    # ===== THE STORY BEGINS =====
    add_section_slide(prs, "The Story", (0, 51, 102))

    add_content_slide(prs, "You Know That Cool Feature?", [
        "OneDrive Photos can transform your photos",
        "into amazing artwork with AI",
        "",
        "Turn yourself into an anime character",
        "Look like a Warhol painting",
        "Become a storybook illustration",
        "",
        "But which style do people love most?"
    ])

    add_big_quote_slide(
        prs,
        "Which style should we show first to users?",
        "The question we needed to answer"
    )

    add_content_slide(prs, "The Old Way: Opinions", [
        "Before, we'd ask around:",
        "",
        '"I think Anime looks cool"',
        '"My favorite is Pop Art"',
        '"Storybook feels warmer to me"',
        "",
        "Everyone had different opinions!",
        "No way to know who was right."
    ])

    add_content_slide(prs, "The New Way: Let AI Be the Judge", [
        "What if we could test this fairly?",
        "",
        "Take real photos",
        "Apply each style",
        "Have AI judges score them",
        "Pick the winner based on data",
        "",
        "That's exactly what we built!"
    ])

    # ===== HOW IT WORKS =====
    add_section_slide(prs, "How It Works", (0, 120, 215))

    add_content_slide(prs, "Think of It Like a Cooking Show", [
        "Same ingredients (3 photos)",
        "Different recipes (3 styles)",
        "Expert judges score each dish",
        "Best overall score wins!",
        "",
        "Fair, consistent, no favoritism"
    ], subtitle="A competition with clear rules")

    add_emoji_content_slide(prs, "The Contestants", [
        "ANIME",
        "Big eyes, bright colors, that cool Japanese cartoon look",
        "",
        "POP ART",
        "Bold colors, dots like comic books, museum-worthy",
        "",
        "STORYBOOK",
        "Soft and dreamy, like a children's book illustration"
    ])

    add_content_slide(prs, "Step 1: Pick the Photos", [
        "We select 3 different photos",
        "",
        "Different people",
        "Different scenes",
        "Real photos from OneDrive",
        "",
        "This makes the test fair"
    ])

    add_content_slide(prs, "Step 2: Transform Each Photo", [
        "Each photo gets all 3 styles applied",
        "",
        "Photo 1 -> Anime, Pop Art, Storybook",
        "Photo 2 -> Anime, Pop Art, Storybook",
        "Photo 3 -> Anime, Pop Art, Storybook",
        "",
        "That's 9 transformed images total!"
    ])

    add_content_slide(prs, "Step 3: The Judges Score", [
        "Two AI judges look at every image:",
        "",
        "Judge #1 (Gemini) asks:",
        '"Does it look right? Any weird glitches?"',
        "",
        "Judge #2 (Claude) asks:",
        '"Would someone want to share this?"'
    ])

    # ===== WHAT WE LOOK FOR =====
    add_section_slide(prs, "What the Judges Look For", (0, 153, 153))

    add_emoji_content_slide(prs, "The Quality Checklist", [
        "Can you still tell who's in the photo?",
        "",
        "Is the style applied everywhere? (no weird patches)",
        "",
        "Does it actually look like that style?",
        "",
        "Are there any weird errors or glitches?",
        "",
        "Would you want to show this to friends?"
    ])

    add_two_column_slide(prs, "Two Different Viewpoints",
        [
            "Checks the technical stuff",
            "",
            "Is it done correctly?",
            "Any broken parts?",
            "Does the style look real?",
            "",
            "Like a quality inspector"
        ],
        [
            "Checks the feeling",
            "",
            "Does it make you smile?",
            "Would you share it?",
            "Is it frame-worthy?",
            "",
            "Like a friend's opinion"
        ],
        "Judge #1: Quality", "Judge #2: Appeal"
    )

    add_content_slide(prs, "Step 4: Pick the Winner", [
        "We combine both judges' scores",
        "",
        "A style needs to be:",
        "Technically good (no glitches)",
        "AND emotionally appealing (makes you go 'wow!')",
        "",
        "The style with best overall score wins!"
    ])

    # ===== WHAT WE GET =====
    add_section_slide(prs, "What We Learn", (102, 51, 153))

    add_emoji_content_slide(prs, "The Results Tell Us", [
        "Which style people will love most",
        "",
        "Which styles need more work",
        "",
        "If quality gets worse over time",
        "",
        "Real data instead of guessing"
    ])

    add_content_slide(prs, "Everything is Saved", [
        "All the original photos",
        "All 9 transformed images",
        "All the scores and reasoning",
        "A final report with the winner",
        "",
        "We can always go back and check!"
    ])

    # ===== WHY IT MATTERS =====
    add_section_slide(prs, "Why This Matters", (0, 153, 76))

    add_two_column_slide(prs, "Before vs After",
        [
            '"I think this one is better"',
            "Different opinions every time",
            "No way to track changes",
            "Takes hours to review",
            "Can't prove anything"
        ],
        [
            '"The data shows this is best"',
            "Same test every time",
            "Track quality over months",
            "Done in minutes",
            "Clear evidence to share"
        ],
        "Guessing", "Knowing"
    )

    add_content_slide(prs, "Better Decisions for Users", [
        "We can confidently say:",
        "",
        '"Put Anime first - users love it most"',
        '"Pop Art needs improvement here..."',
        '"This update made Storybook worse - fix it!"',
        "",
        "Users get better features faster"
    ])

    # ===== OTHER USES =====
    add_section_slide(prs, "What Else Can We Test?", (153, 102, 0))

    add_emoji_content_slide(prs, "This Works for Many Things", [
        "Compare old feature vs new feature",
        "",
        "Test different AI models",
        "",
        "Check if updates break anything",
        "",
        "Compare us vs competitors",
        "",
        "Any time we need to pick a winner!"
    ])

    add_content_slide(prs, "The Pattern", [
        "Take something -> Change it -> Judge it -> Pick best",
        "",
        "Works for:",
        "Photo filters",
        "Video effects",
        "Background removal",
        "Auto-cropping",
        "Any AI feature!"
    ])

    # ===== WRAP UP =====
    add_section_slide(prs, "In Summary", (0, 51, 102))

    add_emoji_content_slide(prs, "What We Built", [
        "An automated way to test AI features",
        "",
        "Fair competition between options",
        "",
        "Two judges with different viewpoints",
        "",
        "Clear winner based on data",
        "",
        "Everything saved for the record"
    ])

    add_big_quote_slide(
        prs,
        "No more guessing.\nNow we know.",
        "Data-driven quality"
    )

    add_content_slide(prs, "Ready to Run!", [
        "The system is set up and tested",
        "",
        "Anime vs Pop Art vs Storybook",
        "3 photos ready to go",
        "Judges standing by",
        "",
        "Let the competition begin!"
    ])

    # Save
    output_path = "Restyle_Benchmark_Story_Simple.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
